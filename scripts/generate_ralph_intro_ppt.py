from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "ralph-overview.pptx"
IMAGE = ROOT / "ralph.webp"


BG = RGBColor(246, 242, 233)
INK = RGBColor(36, 40, 52)
ACCENT = RGBColor(197, 108, 53)
ACCENT_DARK = RGBColor(122, 61, 31)
MUTED = RGBColor(91, 97, 108)
CARD = RGBColor(255, 252, 247)


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_top_bar(slide, label: str):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.35)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.07), Inches(4.5), Inches(0.2))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)


def add_title(slide, title: str, subtitle: str | None = None):
    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.65), Inches(7.3), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = INK

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.67), Inches(1.55), Inches(7.8), Inches(0.9))
        p = sub_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.name = "Aptos"
        run.font.size = Pt(14)
        run.font.color.rgb = MUTED


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(18)
        p.font.color.rgb = INK
        p.space_after = Pt(10)
        p.bullet = True


def add_card(slide, left: float, top: float, width: float, height: float):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = RGBColor(223, 214, 200)
    return card


def add_card_title(slide, text: str, left: float, top: float, width: float):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos Display"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = ACCENT_DARK


def add_small_bullets(slide, items: list[str], left: float, top: float, width: float, height: float):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(15)
        p.font.color.rgb = INK
        p.space_after = Pt(8)
        p.bullet = True


def add_center_callout(slide, text: str):
    box = slide.shapes.add_textbox(Inches(8.25), Inches(1.05), Inches(4.15), Inches(1.7))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos Display"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = ACCENT_DARK


def maybe_add_image(slide):
    if IMAGE.exists():
        try:
            slide.shapes.add_picture(str(IMAGE), Inches(8.7), Inches(1.9), width=Inches(3.7))
            return
        except Exception:
            pass

    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.2), Inches(1.85), Inches(2.7), Inches(2.7)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    box = slide.shapes.add_textbox(Inches(9.5), Inches(2.6), Inches(2.1), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "RALPH"
    run.font.name = "Aptos Display"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_bar(slide, "Ralph oversikt")
    add_title(
        slide,
        "Hva Ralph er",
        "En kort intro til Ralph-løkken og hva som er spesielt med dette repoet.",
    )
    add_center_callout(slide, "Filbasert agentløkke\nfor autonom koding")
    add_bullets(
        slide,
        [
            "Ralph er en gjenopptakbar, PRD-drevet kodingsløkke, ikke bare en enkelt chat-prompt.",
            "Den jobber med én story om gangen i en enagent-flyt med flere iterasjoner.",
            "Tilstand, logger og fremdrift lagres på disk, noe som gjør lange kjøringer enklere å gjenoppta og inspisere.",
        ],
        0.8,
        2.35,
        7.2,
        4.1,
    )
    maybe_add_image(slide)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_bar(slide, "Hvordan Ralph brukes")
    add_title(slide, "Når Ralph passer best", "Det er mest nyttig når du vil ha varig fremdrift over flere steg.")
    add_card(slide, 0.7, 1.8, 5.8, 4.7)
    add_card_title(slide, "Gode bruksområder", 1.0, 2.05, 4.2)
    add_small_bullets(
        slide,
        [
            "Flerstegs repoendringer der arbeidet bør overleve på tvers av iterasjoner.",
            "PRD- eller planstyrt implementering med synlig fremdrift på disk.",
            "Benchmarking, iterativ opprydding og agentlokker som har nytte av lokale logger.",
        ],
        1.0,
        2.45,
        4.9,
        3.5,
    )
    add_card(slide, 6.85, 1.8, 5.8, 4.7)
    add_card_title(slide, "Mindre egnet", 7.15, 2.05, 3.0)
    add_small_bullets(
        slide,
        [
            "Små endringer i én fil der direkte agentarbeid er raskere enn å lage en PRD.",
            "Situasjoner uten lokal verifisering eller reelt tilbakepress fra repoet.",
            "Oppgaver som ikke har nytte av gjenopptakbarhet, logger eller avgrensede iterasjoner.",
        ],
        7.15,
        2.45,
        4.9,
        3.5,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_bar(slide, "Denne repoen vs vanilla Ralph")
    add_title(slide, "Hvordan denne repoen skiller seg ut", "Denne forken er bevisst ikke standard Ralph.")
    add_bullets(
        slide,
        [
            "Windows-først Codex-runner, inkludert SDK-basert supervisjon og roligere håndtering av hjelpeprosesser.",
            "Fersk kontekst per iterasjon, pluss avgrenset fremdrifts-, oppskrifts- og strategiminne for å redusere kontekstrot.",
            "Deterministiske benchmark-suiter som smoke, quick, hourly og deep for repeterbar tuning.",
            "Robusthet for lange kjøringer, som heartbeat-utskrift, hang recovery og innebygde lokale verifiseringsløp.",
        ],
        0.9,
        2.0,
        7.7,
        4.5,
    )
    add_card(slide, 8.95, 2.0, 3.5, 3.2)
    add_card_title(slide, "Kort sagt", 9.25, 2.25, 2.5)
    add_small_bullets(
        slide,
        [
            "Vanilla Ralph er grunnideen for løkken.",
            "Denne repoen er tunet for Codex, Windows, benchmarking og lengre ubetjente kjøringer.",
        ],
        9.25,
        2.7,
        2.8,
        1.9,
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
